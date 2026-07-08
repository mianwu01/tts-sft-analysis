#!/usr/bin/env python3
"""Sewon meeting deck (5 slides), standalone. Style follows slides/make_slides_week_0626.py:
10 x 5.625 in, Georgia bold ~20pt titles, Times New Roman ~10pt tables (header bold, right-aligned
numbers), strict B/W, matplotlib grayscale serif charts. All numbers from the verified tables under
analysis_outputs/data_verification/ (see FINAL_REPORT.md)."""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
DV = os.path.dirname(HERE)
FIG = os.path.join(DV, "figures")
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SHADE = RGBColor(0xEC, 0xEC, 0xEC)
BODY_FONT = "Georgia"
TABLE_FONT = "Times New Roman"

plt.rcParams.update({
    "figure.facecolor": "white", "axes.facecolor": "white", "axes.edgecolor": "black",
    "axes.labelcolor": "black", "xtick.color": "black", "ytick.color": "black",
    "text.color": "black", "font.size": 10, "font.family": "serif",
    "font.serif": ["Georgia", "DejaVu Serif"], "axes.spines.top": False, "axes.spines.right": False,
})


def fig_hard_tail():
    f, axes = plt.subplots(1, 2, figsize=(9.4, 3.0))
    ax = axes[0]
    labels = ["indep-16\n(2nd draw)", "indep-32\n(both draws)", "SE loop1\n(cum 32)",
              "SE loop2\n(cum 48)", "SE union\n(≤loop2)"]
    vals = [1, 1, 38, 51, 58]
    hat = ["", "", "//", "//", "//"]
    bars = ax.bar(labels, vals, color="white", edgecolor="black", width=0.62)
    for b, h in zip(bars, hat):
        b.set_hatch(h)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 1.2, str(v), ha="center", fontsize=10)
    ax.set_ylabel("hard-zeros solved (of 244)")
    ax.set_title("ref-244: 0/16 under the reference draw", fontsize=10)
    ax.set_ylim(0, 66)
    ax.tick_params(axis="x", labelsize=8)
    ax2 = axes[1]
    labels2 = ["indep-16", "indep-32", "SE l1 pop\n(cum32)", "SE≤l1 union", "SE≤l2 union\n(cum48)"]
    vals2 = [850, 850, 872, 887, 907]
    bars2 = ax2.bar(labels2, vals2, color="white", edgecolor="black", width=0.62)
    for b, h in zip(bars2, ["", "", "//", "//", "//"]):
        b.set_hatch(h)
    for b, v in zip(bars2, vals2):
        ax2.text(b.get_x() + b.get_width() / 2, v + 1.5, str(v), ha="center", fontsize=9.5)
    ax2.set_ylim(830, 918)
    ax2.set_ylabel("problems solved (/1,093)")
    ax2.set_title("aggregate coverage at matched budget", fontsize=10)
    ax2.tick_params(axis="x", labelsize=8)
    f.tight_layout()
    p = os.path.join(FIG, "sewon_hardtail.png")
    f.savefig(p, dpi=200)
    plt.close(f)
    return p


HT = fig_hard_tail()

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]
LEFT, CW = 0.38, 9.24


def add_slide(title=None, title_size=20):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(LEFT), Inches(0.23), Inches(CW), Inches(0.55))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = BLACK
        r.font.name = BODY_FONT
    return s


def add_text(slide, x, y, w, h, lines, size=11, leading=1.12):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = leading
        txt, sz, bold = line if isinstance(line, tuple) else (line, size, False)
        r = para.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = BLACK
        r.font.name = BODY_FONT
    return tb


def add_table(slide, rows, x, y, w, font=10, col_widths=None, row_h=0.34, highlight=None):
    highlight = highlight or set()
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(row_h * nr))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(w * cw / total)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SHADE if i in highlight else WHITE
            cell.margin_top = cell.margin_bottom = Pt(1)
            cell.margin_left = cell.margin_right = Pt(4)
            tf = cell.text_frame
            tf.word_wrap = False
            for k, ln in enumerate(str(val).split("\n")):
                para = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                r = para.add_run()
                r.text = ln
                r.font.size = Pt(font)
                r.font.color.rgb = BLACK
                r.font.bold = (i == 0) or (i in highlight)
                r.font.name = TABLE_FONT
                if j > 0:
                    para.alignment = PP_ALIGN.RIGHT
    return shape


def pic(slide, path, x, y, w):
    slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


# ---- Slide 1: main question ----
s = add_slide()
add_text(s, LEFT, 0.6, CW, 4.6, [
    ("Can self-evolution move beyond the initial direct-sampling frontier?", 24, True),
    ("Qwen3-4B · verifier-free SqueezeEvolve (no tests in generation) · 1,093 non-saturated "
     "codeforces · oracle = full hidden tests", 12, False),
    ("", 8, False),
    ("Three pass@K quantities we keep separate:", 13, True),
    ("1.  current/evolving-population pass@K — oracle of the CURRENT candidate set; selection can "
     "never exceed it; we make NO claim against it.", 12, False),
    ("2.  initial direct-sampling pass@K — the base model's first K samples; THIS is the frontier "
     "we test.", 12, False),
    ("3.  compute-matched direct-sampling frontier pass@(N×T) — the same total budget spent on "
     "naive sampling; the “more than just sampling?” baseline.", 12, False),
    ("", 8, False),
    ("Claim under test: evolution reaches solutions that compute-matched independent sampling "
     "essentially never reaches; gains sit at the difficulty boundary; we fold them back into the "
     "same model (self-evolution, no teacher).", 12.5, True),
])

# ---- Slide 2: frontier movement ----
s = add_slide("Frontier movement — evolution vs compute-matched sampling (verified)")
pic(s, HT, 0.45, 0.85, 9.1)
add_text(s, LEFT, 4.05, CW, 1.4, [
    ("A 2nd independent 16-draw adds ~NOTHING (850→850 coverage; 1/244 hard-zeros; the two draws "
     "agree on 243/244) — independent sampling saturates by ~16 samples. SE at the SAME 32-sample "
     "budget solves 38 hard-zeros (51 by loop2, 58 = 23.8% cumulative).", 11.5, True),
    ("Honest caveats: pop-8 cross-check at matched-16: SE 1,384 vs BoN 1,373 (+28/−17) — modest "
     "aggregate edge; independent-48 baseline missing (loop2 not compute-matched yet); partial loop3 "
     "= net −1 (saturating); 7 loop1 solves eroded at loop2 (union keeps them).", 10, False),
])

# ---- Slide 3: boundary buckets ----
s = add_slide("Where gains concentrate — boundary buckets & stopping rules")
pic(s, os.path.join(FIG, "bucket_heatmap.png"), 0.35, 0.85, 5.4)
rows = [
    ["stopping rule (offline sim)", "saved", "frontier kept", "coverage"],
    ["full 2-loop run", "0%", "58/58", "907"],
    ["freeze 9-15 bucket only", "48%", "58/58", "907"],
    ["pure boundary (drop 0-bucket too)", "70%", "1/58", "850"],
    ["stop-when-dry (no new uniques)", "11%", "38/58", "887"],
]
add_table(s, rows, 6.0, 1.0, 3.7, font=9.5, col_widths=[2.0, 0.7, 1.0, 0.9], row_h=0.36,
          highlight={2})
add_text(s, 6.0, 3.0, 3.7, 2.2, [
    ("Boundary buckets (1, 2-4) gain most per loop; 9-15 plateaus at loop1 and its UNIQUE density "
     "falls at loop2 (duplicates).", 10.5, False),
    ("The 0-bucket is thin but carries THE ENTIRE frontier — curricula must keep hard-zero "
     "exploration.", 10.5, True),
    ("(Descriptive simulation on collected loops; no curriculum experiment run yet.)", 9, False),
])

# ---- Slide 4: fold-back into weights ----
s = add_slide("Folding back into weights — SFT pilots (held-out LCBv6, raw generation)")
rows = [
    ["arm (LoRA, empty-think/code-only)", "data", "pass@1", "pass@4", "pass@16", "cap-hit"],
    ["base (untrained)", "—", "0.413", "0.482", "0.527", "44%"],
    ["loop0 · best-per-problem", "804", "0.436", "0.534", "0.595", "24%"],
    ["loop1 (SE) · best-per-problem", "872", "0.437", "0.528", "0.588", "26%"],
    ["loop0 · all-unique", "3,889", "0.410", "0.504", "0.565", "26%"],
    ["loop1 (SE) · all-unique", "9,025", "0.397", "0.494", "0.557", "19%"],
    ["real-CoT · best", "371", "0.203", "0.260", "0.290", "77%"],
]
add_table(s, rows, LEFT, 0.9, CW, font=10, col_widths=[2.6, 0.9, 0.8, 0.8, 0.8, 0.7],
          row_h=0.34, highlight={2, 3})
add_text(s, LEFT, 3.45, CW, 2.0, [
    ("Durable (weight-level): oracle-verified best-per-problem beats base at EVERY k; "
     "SE-source ≈ independent-source at matched selection (Δ 0.001) in this small "
     "cross-platform pilot — SE's 68 extra covered problems can't show up here by construction.", 11, True),
    ("Selection vs durable, kept separate: test-time selection is worth up to +0.10–0.24 pass@1 "
     "on codeforces pools (needs a strong judge: gpt-oss V1 0.80/0.81 ≈ oracle 0.83/0.86; the 4B "
     "can't verify itself — SVD below random). Durable SFT lift so far: +0.023.", 10.5, False),
    ("Negatives verified: all-unique dumping ≤ base @pass@1; real-CoT collapses (78% hit the "
     "40,960-token cap); final-only 0.131. Next: boundary-set RL (Setlur: verifier-free trace-cloning "
     "is suboptimal).", 10, False),
])

# ---- Slide 5: questions ----
s = add_slide("Questions for Sewon")
add_text(s, LEFT, 0.95, CW, 4.4, [
    ("1.  Boundary curriculum: informative band carries the RL gradient, but the 0-bucket carries "
     "the frontier (drop it → lose 57/58 solves). Principled budget split between boundary "
     "exploitation and hard-zero exploration?", 12, False),
    ("2.  Stopping rule: freeze problems at ≥9/16 estimated density? Test-free version would use "
     "gpt-oss SVD density estimates (r≈0.65, over-claim ≤3) — defensible?", 12, False),
    ("3.  SFT vs RL for fold-back: given the SE-vs-independent SFT null at matched selection, is "
     "boundary-set GRPO the cleaner test of evolution-data value?", 12, False),
    ("4.  Final-answer vs evolution-chain EFT: manifests ready (final-correct 9.1k · chain 136 · "
     "pref-pairs 2k), but the SE client logs no parent→child lineage — is per-loop state-chain "
     "supervision a valid proxy, or re-run generation with lineage logging?", 12, False),
    ("5.  Weak self-verifier vs strong judge: 4B self-verification is below random (AUC≈0.41); a "
     "120B judge ≈ oracle. Is judge-in-the-loop honest for a self-evolution story if disclosed?", 12, False),
    ("6.  Sample-efficient RL when rollouts are expensive: beyond compute-matched frontier movement, "
     "what evidence would convince an RL audience that evolution+verification extracts more signal "
     "per rollout?", 12, False),
])

out = os.path.join(HERE, "sewon_meeting_slides.pptx")
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
