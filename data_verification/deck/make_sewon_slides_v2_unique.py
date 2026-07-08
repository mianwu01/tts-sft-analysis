#!/usr/bin/env python3
"""Sewon meeting deck v2 — identical to make_sewon_slides.py EXCEPT slide 3 (the bucket /
stopping-rule slide) is revised to show RAW correct-candidate density vs CANONICAL unique-correct
density side by side (from the unique_correct_audit), with the stopping-rule table on the right.
All other slides are reproduced verbatim. Saves to a NEW file; does not touch the old deck.

Style follows make_sewon_slides.py / slides/make_slides_week_0626.py: 10x5.625 in, Georgia bold
titles, Times New Roman tables, strict B/W. Slide-3 heatmaps are the pre-rendered audit PNGs."""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
DV = os.path.dirname(HERE)
FIG = os.path.join(DV, "figures")
UCA_FIG = os.path.join(DV, "unique_correct_audit", "figures")  # audit heatmaps
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SHADE = RGBColor(0xEC, 0xEC, 0xEC)
BODY_FONT = "Georgia"
TABLE_FONT = "Times New Roman"

plt.rcParams.update({
    "figure.facecolor": "white", "axes.facecolor": "white", "axes.edgecolor": "black",
    "axes.labelcolor": "black", "xtick.color": "black", "ytick.color": "black",
    "text.color": "black", "font.size": 10, "font.family": "serif",
    "font.serif": ["Georgia", "DejaVu Serif"], "axes.spines.top": False, "axes.spines.right": False,
})


def fig_hard_tail():
    f, axes = plt.subplots(1, 2, figsize=(9.4, 3.0))
    ax = axes[0]
    labels = ["indep-16\n(2nd draw)", "indep-32\n(both draws)", "SE loop1\n(cum 32)",
              "SE loop2\n(cum 48)", "SE union\n(≤loop2)"]
    vals = [1, 1, 38, 51, 58]
    hat = ["", "", "//", "//", "//"]
    bars = ax.bar(labels, vals, color="white", edgecolor="black", width=0.62)
    for b, h in zip(bars, hat):
        b.set_hatch(h)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 1.2, str(v), ha="center", fontsize=10)
    ax.set_ylabel("hard-zeros solved (of 244)")
    ax.set_title("ref-244: 0/16 under the reference draw", fontsize=10)
    ax.set_ylim(0, 66)
    ax.tick_params(axis="x", labelsize=8)
    ax2 = axes[1]
    labels2 = ["indep-16", "indep-32", "SE l1 pop\n(cum32)", "SE≤l1 union", "SE≤l2 union\n(cum48)"]
    vals2 = [850, 850, 872, 887, 907]
    bars2 = ax2.bar(labels2, vals2, color="white", edgecolor="black", width=0.62)
    for b, h in zip(bars2, ["", "", "//", "//", "//"]):
        b.set_hatch(h)
    for b, v in zip(bars2, vals2):
        ax2.text(b.get_x() + b.get_width() / 2, v + 1.5, str(v), ha="center", fontsize=9.5)
    ax2.set_ylim(830, 918)
    ax2.set_ylabel("problems solved (/1,093)")
    ax2.set_title("aggregate coverage at matched budget", fontsize=10)
    ax2.tick_params(axis="x", labelsize=8)
    f.tight_layout()
    p = os.path.join(FIG, "sewon_hardtail.png")
    f.savefig(p, dpi=200)
    plt.close(f)
    return p


HT = fig_hard_tail()

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]
LEFT, CW = 0.38, 9.24


def add_slide(title=None, title_size=20):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(LEFT), Inches(0.23), Inches(CW), Inches(0.55))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = BLACK
        r.font.name = BODY_FONT
    return s


def add_text(slide, x, y, w, h, lines, size=11, leading=1.12):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = leading
        txt, sz, bold = line if isinstance(line, tuple) else (line, size, False)
        r = para.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = BLACK
        r.font.name = BODY_FONT
    return tb


def add_table(slide, rows, x, y, w, font=10, col_widths=None, row_h=0.34, highlight=None,
              wrap=False, align_right=True, row_heights=None):
    highlight = highlight or set()
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(row_h * nr))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(w * cw / total)
    if row_heights:
        for i, rh in enumerate(row_heights):
            tbl.rows[i].height = Inches(rh)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SHADE if i in highlight else WHITE
            cell.margin_top = cell.margin_bottom = Pt(1)
            cell.margin_left = cell.margin_right = Pt(4)
            tf = cell.text_frame
            tf.word_wrap = wrap
            for k, ln in enumerate(str(val).split("\n")):
                para = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                r = para.add_run()
                r.text = ln
                r.font.size = Pt(font)
                r.font.color.rgb = BLACK
                r.font.bold = (i == 0) or (i in highlight)
                r.font.name = TABLE_FONT
                if align_right and j > 0:
                    para.alignment = PP_ALIGN.RIGHT
    return shape


def pic(slide, path, x, y, w):
    slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


# ---- Slide 1: main question (verbatim) ----
s = add_slide()
add_text(s, LEFT, 0.6, CW, 4.6, [
    ("Can self-evolution move beyond the initial direct-sampling frontier?", 24, True),
    ("Qwen3-4B · verifier-free SqueezeEvolve (no tests in generation) · 1,093 non-saturated "
     "codeforces · oracle = full hidden tests", 12, False),
    ("", 8, False),
    ("Three pass@K quantities we keep separate:", 13, True),
    ("1.  current/evolving-population pass@K — oracle of the CURRENT candidate set; selection can "
     "never exceed it; we make NO claim against it.", 12, False),
    ("2.  initial direct-sampling pass@K — the base model's first K samples; THIS is the frontier "
     "we test.", 12, False),
    ("3.  compute-matched direct-sampling frontier pass@(N×T) — the same total budget spent on "
     "naive sampling; the “more than just sampling?” baseline.", 12, False),
    ("", 8, False),
    ("Claim under test: evolution reaches solutions that compute-matched independent sampling "
     "essentially never reaches; gains sit at the difficulty boundary; we fold them back into the "
     "same model (self-evolution, no teacher).", 12.5, True),
])

# ---- Slide 2: frontier movement (verbatim) ----
s = add_slide("Frontier movement — evolution vs compute-matched sampling (verified)")
pic(s, HT, 0.45, 0.85, 9.1)
add_text(s, LEFT, 4.05, CW, 1.4, [
    ("A 2nd independent 16-draw adds ~NOTHING (850→850 coverage; 1/244 hard-zeros; the two draws "
     "agree on 243/244) — independent sampling saturates by ~16 samples. SE at the SAME 32-sample "
     "budget solves 38 hard-zeros (51 by loop2, 58 = 23.8% cumulative).", 11.5, True),
    ("Honest caveats: pop-8 cross-check at matched-16: SE 1,384 vs BoN 1,373 (+28/−17) — modest "
     "aggregate edge; independent-48 baseline missing (loop2 not compute-matched yet); partial loop3 "
     "= net −1 (saturating); 7 loop1 solves eroded at loop2 (union keeps them).", 10, False),
])

# ---- Slide 3 (REVISED): raw density vs unique-correct diversity ----
s = add_slide("Where gains concentrate — raw density vs unique-correct diversity")
# left: raw heatmap  |  center: canonical unique-correct heatmap  |  right: stopping-rule table
HMW = 3.62
add_text(s, 0.12, 0.80, HMW, 0.3, [("RAW correct-candidate density (current slide)", 9, True)], leading=1.0)
pic(s, os.path.join(UCA_FIG, "raw_correct_heatmap.png"), 0.12, 1.05, HMW)
add_text(s, 3.86, 0.80, HMW, 0.3, [("CANONICAL unique-correct density (deduplicated)", 9, True)], leading=1.0)
pic(s, os.path.join(UCA_FIG, "canonical_unique_correct_heatmap.png"), 3.86, 1.05, HMW)

add_text(s, 7.60, 0.80, 2.30, 0.3, [("Stopping-rule sim (offline)", 9, True)], leading=1.0)
rows = [
    ["stopping rule", "saved", "kept", "cov"],
    ["full 2-loop run", "0%", "58/58", "907"],
    ["freeze 9–15 only", "48%", "58/58", "907"],
    ["pure boundary\n(drop 0-bucket)", "70%", "1/58", "850"],
    ["stop-when-dry", "11%", "38/58", "887"],
]
add_table(s, rows, 7.58, 1.05, 2.32, font=8, col_widths=[2.05, 0.7, 0.85, 0.7], row_h=0.40,
          highlight={2})
add_text(s, 7.58, 3.35, 2.32, 2.0, [
    ("Freeze the near-saturated 9–15 bucket after loop0 → 48% of evolution compute saved, "
     "ZERO frontier lost (its extra “gains” are duplicates).", 8.5, False),
], leading=1.02)

# caveat + interpretation under the two heatmaps
add_text(s, 0.12, 3.28, 7.32, 2.3, [
    ("Raw density ≠ unique diversity. Dedup confirms hard/boundary gains are genuine, while "
     "easy-bucket gains are mostly copy amplification.", 11, True),
    ("Duplication grows modestly 1.10→1.25→1.33. Easy 9–15 raw rises but canonical-unique falls "
     "11.6→11.5→10.6. Hard/boundary buckets gain genuine unique-correct programs. 0-bucket frontier "
     "is robust: 57 ever-solved; not copy amplification.", 10.5, False),
    ("(Oracle full-test dedup; canonical = comments/whitespace stripped, names & structure kept — "
     "conservative, can only over-count uniqueness. Source: unique_correct_audit.)", 8.5, False),
], leading=1.08)

# ---- Slide 4: fold-back into weights (verbatim) ----
s = add_slide("Folding back into weights — SFT pilots (held-out LCBv6, raw generation)")
rows = [
    ["arm (LoRA, empty-think/code-only)", "data", "pass@1", "pass@4", "pass@16", "cap-hit"],
    ["base (untrained)", "—", "0.413", "0.482", "0.527", "44%"],
    ["loop0 · best-per-problem", "804", "0.436", "0.534", "0.595", "24%"],
    ["loop1 (SE) · best-per-problem", "872", "0.437", "0.528", "0.588", "26%"],
    ["loop0 · all-unique", "3,889", "0.410", "0.504", "0.565", "26%"],
    ["loop1 (SE) · all-unique", "9,025", "0.397", "0.494", "0.557", "19%"],
    ["real-CoT · best", "371", "0.203", "0.260", "0.290", "77%"],
]
add_table(s, rows, LEFT, 0.9, CW, font=10, col_widths=[2.6, 0.9, 0.8, 0.8, 0.8, 0.7],
          row_h=0.34, highlight={2, 3})
add_text(s, LEFT, 3.45, CW, 2.0, [
    ("Durable (weight-level): oracle-verified best-per-problem beats base at EVERY k; "
     "SE-source ≈ independent-source at matched selection (Δ 0.001) in this small "
     "cross-platform pilot — SE's 68 extra covered problems can't show up here by construction.", 11, True),
    ("Selection vs durable, kept separate: test-time selection is worth up to +0.10–0.24 pass@1 "
     "on codeforces pools (needs a strong judge: gpt-oss V1 0.80/0.81 ≈ oracle 0.83/0.86; the 4B "
     "can't verify itself — SVD below random). Durable SFT lift so far: +0.023.", 10.5, False),
    ("Negatives verified: all-unique dumping ≤ base @pass@1; real-CoT collapses (78% hit the "
     "40,960-token cap); final-only 0.131. Next: boundary-set RL (Setlur: verifier-free trace-cloning "
     "is suboptimal).", 10, False),
])

# ---- Slide 5 (NEW): EFT related work ----
s = add_slide("EFT: process-level fold-back works, but via a stronger teacher")
rows = [
    ["EFT evidence", "What they found", "Relevance to us"],
    ["Teacher-generated\ntrajectories",
     "OpenEvolve + Qwen3.5-397B-A17B teacher mutation operator; students are 2B–9B Finch models",
     "Useful but teacher-dependent; not frontier-compatible"],
    ["Held-out gain",
     "+10.22% avg over 22 held-out tasks",
     "Evolution trajectories can become reusable capability"],
    ["Many-task scaling",
     "15→355 training tasks gives +14.1% held-out improvement",
     "Scale / task diversity likely matters"],
    ["Filtered positives\nmatter",
     "Imp-only helps; Imp+NC+Reg can hurt",
     "Matches our all-unique / CoT negative results"],
    ["Test-time RL\nafter EFT",
     "Finch + test-time RL reaches SOTA on two circle-packing tasks and beats base on Erdős",
     "Motivates boundary-set RL after fold-back"],
]
add_table(s, rows, LEFT, 0.9, CW, font=10.5, col_widths=[1.7, 4.0, 2.7], row_h=0.62,
          row_heights=[0.34, 0.68, 0.5, 0.5, 0.5, 0.68], wrap=True, align_right=False, highlight=set())
add_text(s, LEFT, 4.35, CW, 1.15, [
    ("Takeaway: EFT shows evolution-process data can improve held-out discovery ability — but it "
     "relies on a much stronger (~400B) teacher to generate the trajectories.", 12, True),
    ("Our target: same-model self-evolution → verification → fold back into the SAME model "
     "(no teacher).", 12, True),
])

# ---- Slide 6: questions (verbatim) ----
s = add_slide("Questions for Sewon")
add_text(s, LEFT, 0.95, CW, 4.4, [
    ("1.  Boundary curriculum: informative band carries the RL gradient, but the 0-bucket carries "
     "the frontier (drop it → lose 57/58 solves). Principled budget split between boundary "
     "exploitation and hard-zero exploration?", 12, False),
    ("2.  Stopping rule: freeze problems at ≥9/16 estimated density? Test-free version would use "
     "gpt-oss SVD density estimates (r≈0.65, over-claim ≤3) — defensible?", 12, False),
    ("3.  SFT vs RL for fold-back: given the SE-vs-independent SFT null at matched selection, is "
     "boundary-set GRPO the cleaner test of evolution-data value?", 12, False),
    ("4.  Final-answer vs evolution-chain EFT: manifests ready (final-correct 9.1k · chain 136 · "
     "pref-pairs 2k), but the SE client logs no parent→child lineage — is per-loop state-chain "
     "supervision a valid proxy, or re-run generation with lineage logging?", 12, False),
    ("5.  Weak self-verifier vs strong judge: 4B self-verification is below random (AUC≈0.41); a "
     "120B judge ≈ oracle. Is judge-in-the-loop honest for a self-evolution story if disclosed?", 12, False),
    ("6.  Sample-efficient RL when rollouts are expensive: beyond compute-matched frontier movement, "
     "what evidence would convince an RL audience that evolution+verification extracts more signal "
     "per rollout?", 12, False),
])

out = os.path.join(HERE, "sewon_meeting_slides_v2_unique.pptx")
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
